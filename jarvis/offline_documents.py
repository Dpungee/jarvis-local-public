from __future__ import annotations

import html
import hashlib
import json
import os
import stat
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable

from .redaction import redact_secrets


MAX_DOCUMENT_SPEC_BYTES = 1_000_000
MAX_DOCUMENT_OUTPUT_BYTES = 64 * 1024 * 1024
MAX_DOCUMENT_SECTIONS = 100
MAX_SECTION_ITEMS = 500
SUPPORTED_DOCUMENT_TYPES = ("pptx", "docx", "xlsx", "pdf")


@dataclass(frozen=True)
class DocumentSection:
    title: str
    paragraphs: tuple[str, ...] = ()
    bullets: tuple[str, ...] = ()
    rows: tuple[tuple[str, ...], ...] = ()


@dataclass(frozen=True)
class DocumentSpec:
    title: str
    subtitle: str = ""
    sections: tuple[DocumentSection, ...] = ()
    sheet_name: str = "Document"
    sheet_rows: tuple[tuple[str, ...], ...] = ()


def _safe_text(value: Any, limit: int = 10_000) -> str:
    text = redact_secrets(str(value)).replace("\x00", "").strip()
    if len(text) > limit:
        return text[:limit] + "…"
    return text


def _ordinary_file(path: Path, label: str) -> None:
    details = os.lstat(path)
    attributes = getattr(details, "st_file_attributes", 0)
    if (
        stat.S_ISLNK(details.st_mode)
        or attributes & getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0x400)
        or not stat.S_ISREG(details.st_mode)
    ):
        raise PermissionError(f"{label} must be an ordinary file")


def _ordinary_directory(path: Path, label: str) -> None:
    details = os.lstat(path)
    attributes = getattr(details, "st_file_attributes", 0)
    if (
        stat.S_ISLNK(details.st_mode)
        or attributes & getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0x400)
        or not stat.S_ISDIR(details.st_mode)
    ):
        raise PermissionError(f"{label} must be an ordinary directory")


def _inside_workspace(workspace: Path, value: str | os.PathLike[str]) -> Path:
    root = Path(workspace).resolve(strict=True)
    _ordinary_directory(root, "Document workspace")
    raw = Path(value)
    candidate = raw if raw.is_absolute() else root / raw
    resolved = candidate.resolve(strict=False)
    try:
        resolved.relative_to(root)
    except ValueError as exc:
        raise PermissionError("Document paths must stay inside the workspace") from exc
    return resolved


def _reject_link_components(workspace: Path, value: str | os.PathLike[str]) -> Path:
    root = Path(workspace).resolve(strict=True)
    raw = Path(value)
    lexical = raw if raw.is_absolute() else root / raw
    normalized = Path(os.path.abspath(lexical))
    try:
        relative = normalized.relative_to(root)
    except ValueError as exc:
        raise PermissionError("Document paths must stay inside the workspace") from exc
    current = root
    for part in relative.parts:
        current = current / part
        if not os.path.lexists(current):
            continue
        details = os.lstat(current)
        attributes = getattr(details, "st_file_attributes", 0)
        if stat.S_ISLNK(details.st_mode) or attributes & getattr(
            stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0x400
        ):
            raise PermissionError("Document paths may not traverse links")
    return normalized


def _source_path(workspace: Path, source: str | os.PathLike[str]) -> Path:
    _reject_link_components(workspace, source)
    path = _inside_workspace(workspace, source)
    if not path.exists():
        raise FileNotFoundError(f"Document specification was not found: {source}")
    _ordinary_file(path, "Document specification")
    size = path.stat().st_size
    if size <= 0 or size > MAX_DOCUMENT_SPEC_BYTES:
        raise ValueError(
            f"Document specification must contain 1-{MAX_DOCUMENT_SPEC_BYTES} bytes"
        )
    if path.suffix.casefold() not in {".md", ".markdown", ".json"}:
        raise ValueError("Document specification must be Markdown or JSON")
    return path


def _read_document_source(path: Path) -> str:
    before = os.lstat(path)
    with path.open("rb") as stream:
        opened_before = os.fstat(stream.fileno())
        raw = stream.read(MAX_DOCUMENT_SPEC_BYTES + 1)
        opened_after = os.fstat(stream.fileno())
    after = os.lstat(path)
    identities = {
        (item.st_dev, item.st_ino, item.st_size, item.st_mtime_ns)
        for item in (before, opened_before, opened_after, after)
    }
    if len(identities) != 1:
        raise RuntimeError("Document specification changed while it was being read")
    if not raw or len(raw) > MAX_DOCUMENT_SPEC_BYTES:
        raise ValueError(
            f"Document specification must contain 1-{MAX_DOCUMENT_SPEC_BYTES} bytes"
        )
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise ValueError("Document specification must be UTF-8 text") from exc


def _output_path(
    workspace: Path, output: str | os.PathLike[str], document_type: str
) -> Path:
    lexical = _reject_link_components(workspace, output)
    if os.path.lexists(lexical):
        raise FileExistsError("Document output already exists; choose a new path")
    path = _inside_workspace(workspace, output)
    if path.suffix.casefold() != f".{document_type}":
        raise ValueError(f"Output extension must be .{document_type}")
    if not path.parent.exists():
        raise FileNotFoundError("Document output directory does not exist")
    _ordinary_directory(path.parent, "Document output directory")
    return path


def _rows(value: Any) -> tuple[tuple[str, ...], ...]:
    if not isinstance(value, list) or not value:
        return ()
    bounded = value[:MAX_SECTION_ITEMS]
    if all(isinstance(item, dict) for item in bounded):
        keys: list[str] = []
        for item in bounded:
            for key in item:
                normalized = _safe_text(key, 200)
                if normalized and normalized not in keys:
                    keys.append(normalized)
        keys = keys[:100]
        return (
            tuple(keys),
            *(tuple(_safe_text(item.get(key, ""), 2_000) for key in keys) for item in bounded),
        )
    rows: list[tuple[str, ...]] = []
    for item in bounded:
        values = item if isinstance(item, (list, tuple)) else [item]
        rows.append(tuple(_safe_text(cell, 2_000) for cell in values[:100]))
    return tuple(rows)


def _strings(value: Any) -> tuple[str, ...]:
    values = value if isinstance(value, list) else [value]
    return tuple(
        text
        for text in (_safe_text(item) for item in values[:MAX_SECTION_ITEMS])
        if text
    )


def _parse_json(raw: str, fallback_title: str) -> DocumentSpec:
    try:
        value = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise ValueError("Document JSON is malformed") from exc
    if not isinstance(value, dict):
        raise ValueError("Document JSON must contain one top-level object")
    title = _safe_text(value.get("title") or fallback_title, 500) or fallback_title
    subtitle = _safe_text(value.get("subtitle") or "", 1_000)
    sheet_name = _safe_text(value.get("sheet_name") or "Document", 200) or "Document"
    raw_sections = value.get("sections", [])
    if raw_sections is None:
        raw_sections = []
    if not isinstance(raw_sections, list):
        raise ValueError("Document JSON sections must be a list")
    sections: list[DocumentSection] = []
    for index, item in enumerate(raw_sections[:MAX_DOCUMENT_SECTIONS], 1):
        if isinstance(item, str):
            sections.append(DocumentSection(f"Section {index}", (_safe_text(item),)))
            continue
        if not isinstance(item, dict):
            raise ValueError("Every document section must be text or an object")
        body = item.get("paragraphs", item.get("body", []))
        sections.append(
            DocumentSection(
                _safe_text(item.get("title") or f"Section {index}", 500),
                _strings(body) if body not in (None, "") else (),
                _strings(item.get("bullets", [])) if item.get("bullets") else (),
                _rows(item.get("rows", item.get("table", []))),
            )
        )
    top_rows = _rows(value.get("rows", value.get("data", [])))
    if top_rows:
        sections.append(DocumentSection("Data", rows=top_rows))
    sheet_rows = top_rows
    if not sheet_rows:
        row_sections = [
            section for section in sections
            if section.rows and not section.paragraphs and not section.bullets
        ]
        if len(row_sections) == 1:
            sheet_rows = row_sections[0].rows
    return DocumentSpec(
        title,
        subtitle,
        tuple(sections[:MAX_DOCUMENT_SECTIONS]),
        sheet_name,
        sheet_rows,
    )


def _parse_markdown(raw: str, fallback_title: str) -> DocumentSpec:
    title = fallback_title
    subtitle = ""
    found_title = False
    sheet_name = "Document"
    sheet_rows: list[tuple[str, ...]] = []
    sections: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    for original in raw.splitlines():
        line = _safe_text(original)
        if not line:
            continue
        if line.startswith("# ") and not found_title:
            title = _safe_text(line[2:], 500) or fallback_title
            found_title = True
            continue
        if line.casefold().startswith("sheet:"):
            sheet_name = _safe_text(line.split(":", 1)[1], 200) or "Document"
            continue
        if line.startswith("## "):
            if len(sections) >= MAX_DOCUMENT_SECTIONS:
                continue
            current = {
                "title": _safe_text(line[3:], 500) or f"Section {len(sections) + 1}",
                "paragraphs": [],
                "bullets": [],
            }
            sections.append(current)
            continue
        if current is None:
            if line.startswith("|") and line.endswith("|"):
                values = tuple(
                    _safe_text(value, 2_000)
                    for value in line.strip("|").split("|")[:100]
                )
                if values and not all(
                    value and set(value.replace(":", "")) <= {"-"}
                    for value in values
                ):
                    sheet_rows.append(values)
                continue
            if not subtitle:
                subtitle = line
            else:
                current = {"title": "Overview", "paragraphs": [line], "bullets": []}
                sections.append(current)
            continue
        if line.startswith("|") and line.endswith("|"):
            values = tuple(
                _safe_text(value, 2_000)
                for value in line.strip("|").split("|")[:100]
            )
            if values and not all(
                value and set(value.replace(":", "")) <= {"-"}
                for value in values
            ):
                current.setdefault("rows", []).append(values)
                sheet_rows.append(values)
            continue
        target = "bullets" if line.startswith(("- ", "* ", "+ ")) else "paragraphs"
        text = line[2:] if target == "bullets" else line
        if len(current[target]) < MAX_SECTION_ITEMS:
            current[target].append(text)
    parsed = tuple(
        DocumentSection(
            item["title"], tuple(item["paragraphs"]), tuple(item["bullets"]),
            tuple(item.get("rows", ())),
        )
        for item in sections
    )
    return DocumentSpec(
        _safe_text(title, 500),
        _safe_text(subtitle, 1_000),
        parsed,
        sheet_name,
        tuple(sheet_rows[:MAX_SECTION_ITEMS]),
    )


def load_document_spec(path: Path) -> DocumentSpec:
    raw = _read_document_source(path)
    fallback = _safe_text(path.stem.replace("-", " ").replace("_", " ").title(), 500)
    if path.suffix.casefold() == ".json":
        return _parse_json(raw, fallback)
    return _parse_markdown(raw, fallback)


def _document_sections(spec: DocumentSpec) -> tuple[DocumentSection, ...]:
    if spec.sections:
        return spec.sections
    body = (spec.subtitle,) if spec.subtitle else ("Generated from the supplied specification.",)
    return (DocumentSection("Overview", paragraphs=body),)


def _build_pptx(spec: DocumentSpec, output: Path) -> None:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Document libraries are unavailable; install .[documents]") from exc
    deck = Presentation()
    title_slide = deck.slides.add_slide(deck.slide_layouts[0])
    title_slide.shapes.title.text = spec.title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = spec.subtitle
    for section in _document_sections(spec):
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = section.title
        frame = slide.placeholders[1].text_frame
        lines = [*section.paragraphs, *section.bullets]
        if section.rows:
            lines.extend(" | ".join(row) for row in section.rows[:20])
        lines = lines[:20] or [""]
        frame.text = lines[0]
        for line in lines[1:]:
            paragraph = frame.add_paragraph()
            paragraph.text = line
    deck.save(output)


def _build_docx(spec: DocumentSpec, output: Path) -> None:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("Document libraries are unavailable; install .[documents]") from exc
    document = Document()
    document.add_heading(spec.title, level=0)
    if spec.subtitle:
        document.add_paragraph(spec.subtitle)
    for section in _document_sections(spec):
        document.add_heading(section.title, level=1)
        for paragraph in section.paragraphs:
            document.add_paragraph(paragraph)
        for bullet in section.bullets:
            document.add_paragraph(bullet, style="List Bullet")
        if section.rows:
            width = max((len(row) for row in section.rows), default=1)
            table = document.add_table(rows=0, cols=width)
            table.style = "Light Grid Accent 1"
            for values in section.rows:
                cells = table.add_row().cells
                for index, value in enumerate(values):
                    cells[index].text = value
    document.save(output)


def _safe_cell(value: str) -> str:
    return "'" + value if value.startswith(("=", "+", "-", "@")) else value


def _safe_sheet_name(value: str) -> str:
    cleaned = "".join("_" if character in "[]:*?/\\" else character for character in value)
    cleaned = cleaned.strip().strip("'")[:31]
    return cleaned or "Document"


def _build_xlsx(spec: DocumentSpec, output: Path) -> None:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
    except ImportError as exc:
        raise RuntimeError("Document libraries are unavailable; install .[documents]") from exc
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = _safe_sheet_name(spec.sheet_name)
    if spec.sheet_rows:
        for row in spec.sheet_rows:
            sheet.append([_safe_cell(value) for value in row])
        if sheet.max_row:
            for cell in sheet[1]:
                cell.font = Font(bold=True)
                cell.fill = PatternFill("solid", fgColor="D9EAF7")
        sheet.freeze_panes = "A2" if sheet.max_row > 1 else None
        for column in sheet.columns:
            letter = column[0].column_letter
            width = max((len(str(cell.value or "")) for cell in column), default=8)
            sheet.column_dimensions[letter].width = min(60, max(10, width + 2))
        workbook.save(output)
        return
    sheet.append([_safe_cell(spec.title)])
    sheet["A1"].font = Font(bold=True, size=16)
    if spec.subtitle:
        sheet.append([_safe_cell(spec.subtitle)])
    sheet.append([])
    for section in _document_sections(spec):
        sheet.append([_safe_cell(section.title)])
        heading = sheet.cell(row=sheet.max_row, column=1)
        heading.font = Font(bold=True)
        heading.fill = PatternFill("solid", fgColor="D9EAF7")
        for paragraph in section.paragraphs:
            sheet.append([_safe_cell(paragraph)])
        for bullet in section.bullets:
            sheet.append(["• " + _safe_cell(bullet)])
        for row in section.rows:
            sheet.append([_safe_cell(value) for value in row])
        sheet.append([])
    sheet.freeze_panes = "A2"
    sheet.column_dimensions["A"].width = 42
    workbook.save(output)


def _build_pdf(spec: DocumentSpec, output: Path) -> None:
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
    except ImportError as exc:
        raise RuntimeError("Document libraries are unavailable; install .[documents]") from exc
    styles = getSampleStyleSheet()
    story: list[Any] = [Paragraph(html.escape(spec.title), styles["Title"])]
    if spec.subtitle:
        story.extend((Spacer(1, 8), Paragraph(html.escape(spec.subtitle), styles["BodyText"])))
    for section in _document_sections(spec):
        story.extend((Spacer(1, 14), Paragraph(html.escape(section.title), styles["Heading2"])))
        for paragraph in section.paragraphs:
            story.extend((Paragraph(html.escape(paragraph), styles["BodyText"]), Spacer(1, 6)))
        for bullet in section.bullets:
            story.append(Paragraph("• " + html.escape(bullet), styles["BodyText"]))
        if section.rows:
            width = max((len(row) for row in section.rows), default=1)
            table = Table([
                [html.escape(value) for value in (*row, *("" for _ in range(width - len(row))))]
                for row in section.rows
            ])
            table.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
            ]))
            story.extend((Spacer(1, 8), table))
    SimpleDocTemplate(str(output), pagesize=letter).build(story)


_BUILDERS: dict[str, Callable[[DocumentSpec, Path], None]] = {
    "pptx": _build_pptx,
    "docx": _build_docx,
    "xlsx": _build_xlsx,
    "pdf": _build_pdf,
}


def _verify_output(path: Path, document_type: str) -> int:
    _ordinary_file(path, "Generated document")
    size = path.stat().st_size
    if size < 500 or size > MAX_DOCUMENT_OUTPUT_BYTES:
        raise RuntimeError("Generated document failed its size check")
    if document_type in {"pptx", "docx", "xlsx"}:
        if not zipfile.is_zipfile(path):
            raise RuntimeError("Generated Office document is not a valid package")
        expected = {
            "pptx": "ppt/presentation.xml",
            "docx": "word/document.xml",
            "xlsx": "xl/workbook.xml",
        }[document_type]
        with zipfile.ZipFile(path) as archive:
            if expected not in archive.namelist():
                raise RuntimeError("Generated Office document is incomplete")
    elif not path.read_bytes().startswith(b"%PDF-"):
        raise RuntimeError("Generated PDF has an invalid header")
    return size


def build_offline_document(
    workspace: Path,
    source: str | os.PathLike[str],
    output: str | os.PathLike[str],
    document_type: str,
) -> dict[str, Any]:
    kind = str(document_type).strip().casefold()
    if kind not in SUPPORTED_DOCUMENT_TYPES:
        raise ValueError("Document type must be pptx, docx, xlsx, or pdf")
    source_path = _source_path(workspace, source)
    output_path = _output_path(workspace, output, kind)
    spec = load_document_spec(source_path)
    descriptor, temporary_name = tempfile.mkstemp(
        prefix=".jarvis-document-", suffix=f".{kind}", dir=output_path.parent
    )
    os.close(descriptor)
    temporary = Path(temporary_name)
    try:
        _BUILDERS[kind](spec, temporary)
        size = _verify_output(temporary, kind)
        if os.path.lexists(output_path):
            raise FileExistsError("Document output appeared during generation")
        os.replace(temporary, output_path)
        _ordinary_file(output_path, "Generated document")
        result = {
            "type": kind,
            "path": str(output_path),
            "relative_path": output_path.relative_to(Path(workspace).resolve()).as_posix(),
            "bytes": size,
            "sha256": hashlib.sha256(output_path.read_bytes()).hexdigest(),
            "title": spec.title,
        }
        if kind == "xlsx":
            result.update({
                "sheet_names": [_safe_sheet_name(spec.sheet_name)],
                "rows": len(spec.sheet_rows) if spec.sheet_rows else None,
                "columns": (
                    max((len(row) for row in spec.sheet_rows), default=0)
                    if spec.sheet_rows else None
                ),
            })
        return result
    finally:
        temporary.unlink(missing_ok=True)


def build_document_preview(
    workspace: Path,
    source: str | os.PathLike[str],
    output: str | os.PathLike[str],
) -> dict[str, Any]:
    """Create a self-contained, escaped HTML preview and structural QA report."""
    source_path = _source_path(workspace, source)
    lexical = _reject_link_components(workspace, output)
    if os.path.lexists(lexical):
        raise FileExistsError("Document preview already exists; choose a new path")
    output_path = _inside_workspace(workspace, output)
    if output_path.suffix.casefold() != ".html":
        raise ValueError("Document preview output extension must be .html")
    if not output_path.parent.exists():
        raise FileNotFoundError("Document preview output directory does not exist")
    _ordinary_directory(output_path.parent, "Document preview output directory")

    spec = load_document_spec(source_path)
    sections = _document_sections(spec)
    warnings: list[str] = []
    notes: list[str] = []
    if not spec.subtitle:
        notes.append("subtitle_missing")
    for index, section in enumerate(sections, start=1):
        if not (section.paragraphs or section.bullets or section.rows):
            warnings.append(f"section_{index}_empty")
    body: list[str] = [
        "<!doctype html><html lang=\"en\"><head><meta charset=\"utf-8\">",
        "<meta name=\"viewport\" content=\"width=device-width,initial-scale=1\">",
        f"<title>{html.escape(spec.title)}</title>",
        "<style>body{font:16px/1.55 system-ui,sans-serif;max-width:920px;margin:40px auto;padding:0 24px;color:#17202a}h1{margin-bottom:.2em}.subtitle{color:#566573}section{margin:2rem 0}table{border-collapse:collapse;width:100%;overflow-wrap:anywhere}td{border:1px solid #ccd1d1;padding:.45rem;vertical-align:top}tr:first-child{font-weight:700;background:#f4f6f7}@media(prefers-color-scheme:dark){body{background:#111;color:#eee}.subtitle{color:#bbb}td{border-color:#555}tr:first-child{background:#222}}</style></head><body>",
        f"<h1>{html.escape(spec.title)}</h1>",
    ]
    if spec.subtitle:
        body.append(f'<p class="subtitle">{html.escape(spec.subtitle)}</p>')
    for section in sections:
        body.append(f"<section><h2>{html.escape(section.title)}</h2>")
        body.extend(f"<p>{html.escape(item)}</p>" for item in section.paragraphs)
        if section.bullets:
            body.append("<ul>")
            body.extend(f"<li>{html.escape(item)}</li>" for item in section.bullets)
            body.append("</ul>")
        if section.rows:
            body.append("<table>")
            for row in section.rows:
                body.append("<tr>" + "".join(
                    f"<td>{html.escape(value)}</td>" for value in row
                ) + "</tr>")
            body.append("</table>")
        body.append("</section>")
    body.append("</body></html>")
    encoded = "".join(body).encode("utf-8")
    descriptor, temporary_name = tempfile.mkstemp(
        prefix=".jarvis-preview-", suffix=".html", dir=output_path.parent
    )
    temporary = Path(temporary_name)
    try:
        with os.fdopen(descriptor, "wb") as stream:
            stream.write(encoded)
            stream.flush()
            os.fsync(stream.fileno())
        if os.path.lexists(output_path):
            raise FileExistsError("Document preview appeared during generation")
        os.replace(temporary, output_path)
        _ordinary_file(output_path, "Document preview")
    finally:
        temporary.unlink(missing_ok=True)
    return {
        "path": str(output_path),
        "relative_path": output_path.relative_to(Path(workspace).resolve()).as_posix(),
        "bytes": len(encoded),
        "sha256": hashlib.sha256(encoded).hexdigest(),
        "title": spec.title,
        "section_count": len(sections),
        "qa_passed": not warnings,
        "qa_warnings": warnings,
        "qa_notes": notes,
    }
