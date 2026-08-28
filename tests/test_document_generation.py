from __future__ import annotations

import importlib.util
import builtins
import io
import json
import tempfile
import tomllib
import unittest
import zipfile
from contextlib import redirect_stdout
from dataclasses import replace
from pathlib import Path
from unittest.mock import patch

from jarvis import cli
from jarvis.config import Config
from jarvis.memory import Memory
from jarvis.offline_documents import build_document_preview, build_offline_document
from jarvis.skill_library import read_builtin_skill
from jarvis.tools import ToolBox


ROOT = Path(__file__).resolve().parents[1]


class DocumentGenerationTests(unittest.TestCase):
    def test_html_preview_is_escaped_atomic_and_reports_structural_qa(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            workspace = Path(directory)
            (workspace / "preview.md").write_text(
                "# <Quarterly & Review>\n\n## Results\n<script>alert(1)</script>\n\n"
                "## Empty\n",
                encoding="utf-8",
            )
            result = build_document_preview(workspace, "preview.md", "preview.html")
            content = (workspace / "preview.html").read_text(encoding="utf-8")
            self.assertIn("&lt;Quarterly &amp; Review&gt;", content)
            self.assertIn("&lt;script&gt;alert(1)&lt;/script&gt;", content)
            self.assertNotIn("<script>alert(1)</script>", content)
            self.assertEqual(result["section_count"], 2)
            self.assertFalse(result["qa_passed"])
            self.assertIn("subtitle_missing", result["qa_notes"])
            self.assertIn("section_2_empty", result["qa_warnings"])
            self.assertEqual(list(workspace.glob(".jarvis-preview-*")), [])

            with self.assertRaises(FileExistsError):
                build_document_preview(workspace, "preview.md", "preview.html")
            with self.assertRaises(PermissionError):
                build_document_preview(workspace, "preview.md", "../escape.html")
    def test_document_extra_declares_every_supported_library(self) -> None:
        manifest = tomllib.loads(ROOT.joinpath("pyproject.toml").read_text(encoding="utf-8"))
        requirements = manifest["project"]["optional-dependencies"]["documents"]
        names = {item.split("=", 1)[0].split(">", 1)[0].casefold() for item in requirements}
        self.assertEqual(
            names,
            {
                "python-pptx",
                "python-docx",
                "openpyxl",
                "reportlab",
                "pypdf",
                "pillow",
                "matplotlib",
            },
        )

    def test_document_skill_is_loadable_and_bounded(self) -> None:
        skill = read_builtin_skill("document-generation")
        self.assertLessEqual(len(skill["content"].encode("utf-8")), 32 * 1024)
        self.assertIn("build_document", skill["content"])
        self.assertIn("exact workspace path", skill["content"])

    def test_fresh_setup_installs_only_the_declared_document_extra(self) -> None:
        setup = ROOT.joinpath("setup.ps1").read_text(encoding="utf-8")
        self.assertIn('"--no-input", "--editable", ".[documents]"', setup)
        self.assertNotIn("pip install <name>", setup)

    @unittest.skipUnless(
        all(importlib.util.find_spec(name) for name in ("pptx", "docx", "openpyxl", "reportlab")),
        "document-generation extra is not installed in this test environment",
    )
    def test_declared_libraries_create_real_minimal_documents(self) -> None:
        from docx import Document
        from openpyxl import Workbook
        from pptx import Presentation
        from reportlab.pdfgen import canvas

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            deck = Presentation()
            deck.slides.add_slide(deck.slide_layouts[0]).shapes.title.text = "Verified"
            deck.save(root / "verified.pptx")

            document = Document()
            document.add_heading("Verified", level=0)
            document.save(root / "verified.docx")

            workbook = Workbook()
            workbook.active["A1"] = "Verified"
            workbook.save(root / "verified.xlsx")

            pdf = canvas.Canvas(str(root / "verified.pdf"))
            pdf.drawString(72, 720, "Verified")
            pdf.save()

            for name in ("verified.pptx", "verified.docx", "verified.xlsx", "verified.pdf"):
                with self.subTest(name=name):
                    self.assertGreater(root.joinpath(name).stat().st_size, 500)

    @unittest.skipUnless(
        all(importlib.util.find_spec(name) for name in ("pptx", "docx", "openpyxl", "reportlab")),
        "document-generation extra is not installed in this test environment",
    )
    def test_offline_builder_creates_all_formats_without_a_model(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            workspace = Path(directory)
            source = workspace / "quarterly.md"
            source.write_text(
                "# Quarterly Review\nPrepared for the team\n\n"
                "## Highlights\n- Revenue increased\n- Retention improved\n\n"
                "## Next steps\nShip the verified roadmap.\n",
                encoding="utf-8",
            )
            expected_members = {
                "pptx": "ppt/presentation.xml",
                "docx": "word/document.xml",
                "xlsx": "xl/workbook.xml",
            }
            for kind in ("pptx", "docx", "xlsx", "pdf"):
                with self.subTest(kind=kind):
                    result = build_offline_document(
                        workspace, "quarterly.md", f"quarterly.{kind}", kind
                    )
                    output = workspace / f"quarterly.{kind}"
                    self.assertEqual(result["relative_path"], f"quarterly.{kind}")
                    self.assertGreater(output.stat().st_size, 500)
                    if kind in expected_members:
                        with zipfile.ZipFile(output) as archive:
                            self.assertIn(expected_members[kind], archive.namelist())
                    else:
                        self.assertTrue(output.read_bytes().startswith(b"%PDF-"))

    @unittest.skipUnless(
        importlib.util.find_spec("docx") and importlib.util.find_spec("openpyxl"),
        "document-generation extra is not installed in this test environment",
    )
    def test_offline_builder_redacts_secrets_and_blocks_spreadsheet_formulas(self) -> None:
        from docx import Document
        from openpyxl import load_workbook

        with tempfile.TemporaryDirectory() as directory:
            workspace = Path(directory)
            secret = "sk-proj-" + "A" * 32
            # Deliberate fake canary proves the generated document redacts secrets.
            # codeql[py/clear-text-storage-sensitive-data]
            (workspace / "report.json").write_text(
                json.dumps({
                    "title": "Safe report",
                    "sections": [{
                        "title": "Results",
                        "paragraphs": [f"api_key={secret}"],
                        "rows": [["Name", "Value"], ["Example", "=1+1"]],
                    }],
                }),
                encoding="utf-8",
            )
            build_offline_document(workspace, "report.json", "safe.docx", "docx")
            text = "\n".join(
                paragraph.text for paragraph in Document(workspace / "safe.docx").paragraphs
            )
            self.assertNotIn(secret, text)
            self.assertIn("[REDACTED]", text)

            build_offline_document(workspace, "report.json", "safe.xlsx", "xlsx")
            workbook = load_workbook(workspace / "safe.xlsx", data_only=False)
            values = [cell.value for row in workbook.active.iter_rows() for cell in row]
            self.assertIn("'=1+1", values)
            self.assertNotIn("=1+1", values)

    def test_offline_builder_rejects_escape_overwrite_and_wrong_extension(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            workspace = Path(directory) / "workspace"
            workspace.mkdir()
            (workspace / "spec.md").write_text("# Safe\nBody\n", encoding="utf-8")
            (workspace / "existing.pdf").write_bytes(b"existing")
            with self.assertRaises(PermissionError):
                build_offline_document(workspace, "spec.md", "../escape.pdf", "pdf")
            with self.assertRaises(FileExistsError):
                build_offline_document(workspace, "spec.md", "existing.pdf", "pdf")
            with self.assertRaisesRegex(ValueError, "extension"):
                build_offline_document(workspace, "spec.md", "wrong.docx", "pdf")

    @unittest.skipUnless(
        importlib.util.find_spec("docx"),
        "document-generation extra is not installed in this test environment",
    )
    def test_toolbox_build_document_is_direct_verified_and_cleans_source(self) -> None:
        from docx import Document

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            workspace = root / "workspace"
            data = root / "data"
            workspace.mkdir()
            data.mkdir()
            config = replace(
                Config.load(),
                workspace=workspace,
                data_dir=data,
                vault_dir=None,
                autonomy="autonomous",
                execution_mode="disabled",
                computer_access="disabled",
                external_access="disabled",
            )
            memory = Memory(data / "jarvis.db")
            try:
                toolbox = ToolBox(config, memory)
                result = json.loads(toolbox.execute("build_document", {
                    "path": "brief.docx",
                    "document_type": "docx",
                    "content": "# Speed Brief\n\n## Result\n- Fast\n- Verified\n",
                }))
                self.assertTrue(result["ok"])
                self.assertTrue(result["result"]["verified"])
                self.assertEqual(result["result"]["relative_path"], "brief.docx")
                paragraphs = Document(workspace / "brief.docx").paragraphs
                self.assertIn("Speed Brief", [item.text for item in paragraphs])
                self.assertEqual(
                    list(workspace.glob(".jarvis-document-source-*.md")),
                    [],
                )
            finally:
                memory.close()

    def test_offline_builder_reports_missing_document_extra_clearly(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            workspace = Path(directory)
            (workspace / "spec.md").write_text("# Missing library check\n", encoding="utf-8")
            original_import = builtins.__import__

            def blocked_import(name, *args, **kwargs):
                if name == "pptx":
                    raise ImportError("synthetic missing library")
                return original_import(name, *args, **kwargs)

            with (
                patch("builtins.__import__", side_effect=blocked_import),
                self.assertRaisesRegex(RuntimeError, r"install \.\[documents\]"),
            ):
                build_offline_document(workspace, "spec.md", "missing.pptx", "pptx")
            self.assertFalse((workspace / "missing.pptx").exists())

    @unittest.skipUnless(
        importlib.util.find_spec("docx"),
        "document-generation extra is not installed in this test environment",
    )
    def test_doc_cli_parses_and_builds_without_provider_setup(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            workspace = Path(directory)
            (workspace / "spec.md").write_text("# Offline CLI\nReady.\n", encoding="utf-8")
            config = replace(Config.load(), workspace=workspace)
            args = cli._parser().parse_args(
                ["doc", "--type", "docx", "--from", "spec.md", "result.docx"]
            )
            output = io.StringIO()
            with patch.object(cli.Config, "load", return_value=config), redirect_stdout(output):
                self.assertEqual(cli._run_doc(args), 0)
            self.assertTrue((workspace / "result.docx").is_file())
            self.assertIn("Created DOCX document", output.getvalue())
